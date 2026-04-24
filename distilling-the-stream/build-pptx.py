#!/usr/bin/env python3
"""
Build a PowerPoint file from distilling-the-stream.html.

Pipeline:
  1. Chrome headless renders the HTML to a PDF, one page per slide.
  2. pypdfium2 rasterizes each PDF page to a PNG.
  3. lxml extracts the speaker note (HTML <!-- SPEAKER: ... --> comment) from
     each top-level <div class="slide"> in document order.
  4. python-pptx assembles a 13.333 x 7.5 in (16:9) deck: one full-bleed image
     per slide, speaker notes attached to the notes pane.
"""

from __future__ import annotations

import pathlib
import shutil
import subprocess
import sys
import tempfile

from lxml import etree, html
import pypdfium2 as pdfium
from pptx import Presentation
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
HTML_FILE = HERE / "distilling-the-stream.html"
OUT_PPTX = HERE / "distilling-the-stream.pptx"
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

# 1280x720 at 96dpi -> 13.333 x 7.5 inches. This is the standard 16:9 slide
# size in PowerPoint.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Rasterization scale. 2.0 -> 2560x1440 per slide; plenty crisp at any zoom.
RASTER_SCALE = 2.0


def render_pdf(html_path: pathlib.Path, pdf_path: pathlib.Path) -> None:
    """Run headless Chrome to print the HTML to a PDF at exact slide size."""
    if not pathlib.Path(CHROME).exists():
        sys.exit(f"Chrome not found at {CHROME}")

    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-pdf-header-footer",
        "--hide-scrollbars",
        # wait up to 10s of virtual time so web fonts / images settle
        "--virtual-time-budget=10000",
        f"--print-to-pdf={pdf_path}",
        html_path.as_uri(),
    ]
    print(f"[1/4] rendering PDF -> {pdf_path.name}")
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0 or not pdf_path.exists():
        sys.stderr.write(result.stdout)
        sys.stderr.write(result.stderr)
        sys.exit(f"Chrome failed (exit {result.returncode}).")


def _is_blank(pil_image) -> bool:
    """
    Heuristic: a page is "blank" if almost all pixels cluster at the
    background color. The real slides have text, icons, panels, QR codes —
    they have a lot of pixels far from the background. Blank pages have
    only the grain-texture overlay, which is near-uniform.
    """
    from PIL import Image
    # downscale to cap work; convert to grayscale
    thumb = pil_image.convert("L").resize((160, 90), Image.BILINEAR)
    px = list(thumb.getdata())
    # Background (--bg-deep = #0A0B12) maps to grayscale ~11.
    # Count pixels that differ from background by more than a tolerance.
    non_bg = sum(1 for p in px if abs(p - 11) > 20)
    ratio = non_bg / len(px)
    return ratio < 0.01  # <1% of pixels carry any real content


def rasterize(pdf_path: pathlib.Path, out_dir: pathlib.Path) -> list[pathlib.Path]:
    """
    Rasterize every page of the PDF to a PNG. Silently drops blank pages
    (Chrome occasionally emits a leading or trailing blank from page-break
    interactions). Returns remaining paths in document order.
    """
    print(f"[2/4] rasterizing PDF pages (scale={RASTER_SCALE})")
    pdf = pdfium.PdfDocument(pdf_path)
    paths: list[pathlib.Path] = []
    blanks: list[int] = []
    for i, page in enumerate(pdf, start=1):
        pil_image = page.render(scale=RASTER_SCALE).to_pil()
        if _is_blank(pil_image):
            blanks.append(i)
            continue
        out = out_dir / f"slide-{len(paths)+1:02d}.png"
        pil_image.save(out, "PNG", optimize=True)
        paths.append(out)
    pdf.close()
    if blanks:
        print(f"      dropped blank page(s) at pdf index: {blanks}")
    print(f"      -> {len(paths)} content page(s)")
    return paths


def extract_speaker_notes(html_path: pathlib.Path) -> list[str]:
    """
    Return one speaker-note string per top-level slide, in document order.
    A slide without a SPEAKER comment gets an empty string.
    """
    print("[3/4] extracting speaker notes from HTML")
    tree = html.parse(str(html_path))
    body = tree.getroot().find("body")
    if body is None:
        sys.exit("could not find <body> in HTML")
    slides = [
        el for el in body
        if el.tag == "div" and "slide" in (el.get("class") or "").split()
    ]
    notes: list[str] = []
    for slide in slides:
        speaker = ""
        for c in slide.iter(etree.Comment):
            text = (c.text or "").strip()
            if text.startswith("SPEAKER:"):
                speaker = text[len("SPEAKER:"):].strip()
                break
        notes.append(speaker)
    print(f"      -> {len(notes)} slide(s), {sum(1 for n in notes if n)} with notes")
    return notes


def build_pptx(images: list[pathlib.Path], notes: list[str], out_path: pathlib.Path) -> None:
    """Assemble the PPTX: one full-bleed image per slide, notes attached."""
    print(f"[4/4] assembling PPTX -> {out_path.name}")
    if len(images) != len(notes):
        print(
            f"      warn: {len(images)} images vs {len(notes)} slides — "
            "alignment may be off (check blank-page detection)"
        )

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]  # blank

    for idx, img in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        # Attach speaker note if we have one for this slide.
        note_text = notes[idx] if idx < len(notes) else ""
        if note_text:
            tf = slide.notes_slide.notes_text_frame
            tf.text = note_text
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

    prs.save(out_path)


def main() -> None:
    if not HTML_FILE.exists():
        sys.exit(f"HTML source not found: {HTML_FILE}")

    with tempfile.TemporaryDirectory(prefix="dts-pptx-") as tmp:
        tmp_dir = pathlib.Path(tmp)
        pdf_path = tmp_dir / "deck.pdf"

        render_pdf(HTML_FILE, pdf_path)
        images = rasterize(pdf_path, tmp_dir)
        notes = extract_speaker_notes(HTML_FILE)
        build_pptx(images, notes, OUT_PPTX)

    size_kb = OUT_PPTX.stat().st_size // 1024
    print(f"done: {OUT_PPTX} ({size_kb} KB)")


if __name__ == "__main__":
    main()
